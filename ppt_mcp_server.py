from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
import os

mcp = FastMCP("ppt-server")

BASE_DIR = os.path.abspath("ppt_storage")
ASSETS_DIR = os.path.abspath("assets")
os.makedirs(BASE_DIR, exist_ok=True)

ppt = None
ppt_path = None


@mcp.tool()
def create_presentation(filename: str, topic: str = "default") -> str:
    global ppt, ppt_path
    ppt = Presentation()
    ppt_path = os.path.join(BASE_DIR, filename)
    return "Presentation created"


@mcp.tool()
def add_slide(title: str, content: str, topic: str = "default") -> str:
    global ppt

    slide_layout = ppt.slide_layouts[6]  # blank layout
    slide = ppt.slides.add_slide(slide_layout)

    # 🖼️ Background image
    img_path = os.path.join(ASSETS_DIR, "space.png")
    slide.shapes.add_picture(img_path, 0, 0, width=ppt.slide_width, height=ppt.slide_height)

    # 🧱 Title box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 223, 0)  # Gold

    # 🧱 Content box
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    first = True
    for line in content.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(255, 255, 255)  # White

    return f"Styled slide: {title}"


@mcp.tool()
def save_presentation() -> str:
    global ppt, ppt_path
    ppt.save(ppt_path)
    return f"Saved at {ppt_path}"


if __name__ == "__main__":
    mcp.run()
