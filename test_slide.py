from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
import os

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[6])

img_path = os.path.abspath("assets/space.png")
print("Image exists:", os.path.exists(img_path))
print("Image path:", img_path)

slide.shapes.add_picture(img_path, 0, 0, width=ppt.slide_width, height=ppt.slide_height)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Test Title - Solar System"
run.font.size = Pt(36)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 223, 0)

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(4.5))
tf = content_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "This is test content"
run.font.size = Pt(22)
run.font.color.rgb = RGBColor(255, 255, 255)

out = "ppt_storage/test_output.pptx"
ppt.save(out)
print("Saved to:", os.path.abspath(out))
print("Shapes on slide:", len(slide.shapes))
for s in slide.shapes:
    print(" -", s.shape_type, s.name)
