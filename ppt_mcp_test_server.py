from mcp.server.fastmcp import FastMCP
from pptx import Presentation
import os

mcp = FastMCP("ppt-test-server")

BASE_DIR = os.path.abspath("ppt_storage")
os.makedirs(BASE_DIR, exist_ok=True)

ppt = None
ppt_path = None


@mcp.tool()
def create_presentation(filename: str) -> str:
    global ppt, ppt_path
    ppt = Presentation()
    ppt_path = os.path.join(BASE_DIR, filename)
    return f"Presentation created: {ppt_path}"


@mcp.tool()
def add_slide(title: str, content: str) -> str:
    global ppt

    if ppt is None:
        return "Error: Create presentation first"

    slide_layout = ppt.slide_layouts[1]
    slide = ppt.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = content

    return f"Slide added: {title}"


@mcp.tool()
def save_presentation() -> str:
    global ppt, ppt_path

    if ppt is None:
        return "Error: No presentation to save"

    ppt.save(ppt_path)
    return f"Saved at {ppt_path}"


if __name__ == "__main__":
    mcp.run()